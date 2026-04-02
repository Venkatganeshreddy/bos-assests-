from __future__ import annotations

import csv
import io
import json
import os
import re
from dataclasses import dataclass
from typing import Any

import pandas as pd
import streamlit as st
from config_runtime import get_json_env, get_secret_mapping
from docx import Document as DocxDocument
from google.oauth2.service_account import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from pypdf import PdfReader
from pptx import Presentation

SCOPES = [
    "https://www.googleapis.com/auth/drive.readonly",
    "https://www.googleapis.com/auth/spreadsheets.readonly",
]

FOLDER_MIME = "application/vnd.google-apps.folder"
SHORTCUT_MIME = "application/vnd.google-apps.shortcut"
GOOGLE_DOC_MIME = "application/vnd.google-apps.document"
GOOGLE_SHEET_MIME = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDE_MIME = "application/vnd.google-apps.presentation"

TEXT_FILE_SUFFIXES = {
    ".txt",
    ".md",
    ".csv",
    ".tsv",
    ".json",
    ".yaml",
    ".yml",
    ".html",
    ".htm",
    ".xml",
    ".rtf",
    ".svg",
}

IMAGE_FILE_SUFFIXES = {
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".gif",
    ".bmp",
    ".tif",
    ".tiff",
}

SHEET_META_FIELDS = "properties.title,sheets.properties.title"
MAX_SHEET_ROWS = 401
MAX_INDEX_FILE_BYTES = 8 * 1024 * 1024


@dataclass
class DriveDocument:
    file_id: str
    name: str
    mime_type: str
    path: str
    web_view_link: str
    modified_time: str
    source_kind: str
    text: str


def extract_folder_id(folder_url: str) -> str:
    match = re.search(r"/folders/([a-zA-Z0-9_-]+)", folder_url)
    if match:
        return match.group(1)

    match = re.search(r"[?&]id=([a-zA-Z0-9_-]+)", folder_url)
    if match:
        return match.group(1)

    if re.fullmatch(r"[a-zA-Z0-9_-]{10,}", folder_url.strip()):
        return folder_url.strip()

    raise ValueError("Could not find a Google Drive folder ID in the provided URL.")


def get_service_account_email_from_config() -> str:
    credentials = load_credentials()
    return credentials.service_account_email


def load_credentials() -> Credentials:
    secret_info = get_secret_mapping("gcp_service_account", "google_service_account")
    if secret_info:
        return Credentials.from_service_account_info(secret_info, scopes=SCOPES)

    env_info = get_json_env("GCP_SERVICE_ACCOUNT_JSON", "GOOGLE_SERVICE_ACCOUNT_JSON")
    if env_info:
        return Credentials.from_service_account_info(env_info, scopes=SCOPES)

    credentials_path = os.getenv("GOOGLE_CREDENTIALS_PATH", "credentials.json")
    if not os.path.exists(credentials_path):
        raise FileNotFoundError(
            "Google credentials not found. Set GOOGLE_CREDENTIALS_PATH or add "
            "`gcp_service_account` or `google_service_account` to Streamlit secrets."
        )

    return Credentials.from_service_account_file(credentials_path, scopes=SCOPES)


def build_clients() -> tuple[Any, Any]:
    credentials = load_credentials()
    drive_service = build("drive", "v3", credentials=credentials, cache_discovery=False)
    sheets_service = build("sheets", "v4", credentials=credentials, cache_discovery=False)
    return drive_service, sheets_service


def index_drive_folder(folder_url: str) -> dict[str, Any]:
    return index_drive_folder_with_options(
        folder_url=folder_url,
    )


def index_drive_folder_with_options(
    folder_url: str,
    progress_callback=None,
) -> dict[str, Any]:
    folder_id = extract_folder_id(folder_url)
    credentials = load_credentials()
    drive_service = build("drive", "v3", credentials=credentials, cache_discovery=False)
    sheets_service = build("sheets", "v4", credentials=credentials, cache_discovery=False)

    folder_meta = drive_service.files().get(
        fileId=folder_id,
        fields="id,name,mimeType,webViewLink",
        supportsAllDrives=True,
    ).execute()

    if folder_meta.get("mimeType") != FOLDER_MIME:
        raise ValueError("The provided URL does not point to a Google Drive folder.")

    documents: list[DriveDocument] = []
    skipped: list[dict[str, str]] = []

    items = _walk_folder(
        drive_service=drive_service,
        folder_id=folder_id,
        prefix=folder_meta.get("name", "BOS Assets"),
        visited_folders=set(),
    )

    total_items = len(items)
    _report_progress(progress_callback, 0, total_items, "Scanning BOS assets...")

    for index, item in enumerate(items, start=1):
        size_bytes = _coerce_size_bytes(item.get("size"))
        if size_bytes is not None and size_bytes > MAX_INDEX_FILE_BYTES:
            skipped.append(
                {
                    "name": item["path"],
                    "reason": f"Skipped because file is larger than {MAX_INDEX_FILE_BYTES // (1024 * 1024)} MB.",
                    "link": _item_open_link(item),
                }
            )
            _report_progress(progress_callback, index, total_items, item["path"])
            continue

        extracted = _extract_item_for_index(
            _drive_service=drive_service,
            _sheets_service=sheets_service,
            file_id=item["effectiveId"],
            mime_type=item["effectiveMimeType"],
            name=item["name"],
            path=item["path"],
            web_view_link=item.get("webViewLink", ""),
            modified_time=item.get("modifiedTime", ""),
        )

        if extracted.get("document"):
            documents.append(DriveDocument(**extracted["document"]))
        else:
            skipped.append(extracted["skipped"])

        _report_progress(progress_callback, index, total_items, item["path"])

    return {
        "folder_id": folder_id,
        "folder_name": folder_meta.get("name", "BOS Assets"),
        "folder_link": folder_meta.get("webViewLink", folder_url),
        "service_account_email": credentials.service_account_email,
        "documents": documents,
        "skipped": skipped,
    }


def _report_progress(callback, current: int, total: int, label: str = "") -> None:
    if not callback:
        return
    callback(current, total, label)


def _coerce_size_bytes(raw_size: Any) -> int | None:
    if raw_size in (None, ""):
        return None
    try:
        return int(raw_size)
    except (TypeError, ValueError):
        return None


@st.cache_data(ttl=43200, show_spinner=False)
def _extract_item_for_index(
    _drive_service: Any,
    _sheets_service: Any,
    *,
    file_id: str,
    mime_type: str,
    name: str,
    path: str,
    web_view_link: str,
    modified_time: str,
) -> dict[str, Any]:
    item = {
        "effectiveId": file_id,
        "effectiveMimeType": mime_type,
        "name": name,
        "path": path,
        "webViewLink": web_view_link,
        "modifiedTime": modified_time,
    }

    try:
        text, source_kind = _extract_item_text(
            drive_service=_drive_service,
            sheets_service=_sheets_service,
            item=item,
        )
        normalized = _normalize_extracted_text(text)
        if not normalized:
            return {
                "document": None,
                "skipped": {
                    "name": path,
                    "reason": "The file did not contain readable text.",
                    "link": _item_open_link(item),
                },
            }

        return {
            "document": {
                "file_id": file_id,
                "name": name,
                "mime_type": mime_type,
                "path": path,
                "web_view_link": web_view_link,
                "modified_time": modified_time,
                "source_kind": source_kind,
                "text": normalized,
            },
            "skipped": None,
        }
    except Exception as exc:
        return {
            "document": None,
            "skipped": {
                "name": path,
                "reason": str(exc)[:220],
                "link": _item_open_link(item),
            },
        }


def _walk_folder(
    drive_service: Any,
    folder_id: str,
    prefix: str,
    visited_folders: set[str],
) -> list[dict[str, Any]]:
    if folder_id in visited_folders:
        return []
    visited_folders.add(folder_id)

    items: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = drive_service.files().list(
            q=f"'{folder_id}' in parents and trashed = false",
            pageSize=1000,
            pageToken=page_token,
            includeItemsFromAllDrives=True,
            supportsAllDrives=True,
            fields=(
                "nextPageToken,"
                "files(id,name,mimeType,modifiedTime,size,webViewLink,"
                "shortcutDetails(targetId,targetMimeType))"
            ),
        ).execute()
        items.extend(response.get("files", []))
        page_token = response.get("nextPageToken")
        if not page_token:
            break

    items.sort(
        key=lambda item: (
            _effective_mime_type(item) != FOLDER_MIME,
            item.get("name", "").lower(),
        )
    )

    collected: list[dict[str, Any]] = []
    for item in items:
        effective_id = _effective_id(item)
        effective_mime = _effective_mime_type(item)
        item_path = f"{prefix}/{item.get('name', 'Untitled')}"

        if effective_mime == FOLDER_MIME:
            collected.extend(
                _walk_folder(
                    drive_service=drive_service,
                    folder_id=effective_id,
                    prefix=item_path,
                    visited_folders=visited_folders,
                )
            )
            continue

        collected.append(
            {
                **item,
                "effectiveId": effective_id,
                "effectiveMimeType": effective_mime,
                "path": item_path,
            }
        )

    return collected


def _effective_id(item: dict[str, Any]) -> str:
    details = item.get("shortcutDetails") or {}
    return details.get("targetId") or item["id"]


def _effective_mime_type(item: dict[str, Any]) -> str:
    details = item.get("shortcutDetails") or {}
    return details.get("targetMimeType") or item["mimeType"]


def _item_open_link(item: dict[str, Any]) -> str:
    effective_id = item.get("effectiveId") or _effective_id(item)
    explicit_link = item.get("webViewLink", "")
    if explicit_link:
        return explicit_link
    if effective_id:
        return f"https://drive.google.com/open?id={effective_id}"
    return ""


def _extract_item_text(
    drive_service: Any,
    sheets_service: Any,
    item: dict[str, Any],
) -> tuple[str, str]:
    file_id = item["effectiveId"]
    mime_type = item["effectiveMimeType"]
    name = item["name"]
    lowered_name = name.lower()

    if mime_type == GOOGLE_DOC_MIME:
        return _export_google_workspace_text(drive_service, file_id), "google-doc"

    if mime_type == GOOGLE_SLIDE_MIME:
        return _export_google_workspace_text(drive_service, file_id), "google-slide"

    if mime_type == GOOGLE_SHEET_MIME:
        return _render_google_sheet(sheets_service, file_id), "google-sheet"

    if mime_type == "image/svg+xml" or lowered_name.endswith(".svg"):
        file_bytes = _download_file_bytes(drive_service, file_id)
        return _render_text_like(file_bytes, lowered_name), "svg"

    if mime_type == "application/pdf" or lowered_name.endswith(".pdf"):
        file_bytes = _download_file_bytes(drive_service, file_id)
        rendered_pdf_text = _render_pdf(file_bytes)
        if rendered_pdf_text.strip():
            return rendered_pdf_text, "pdf"
        raise ValueError("Scanned or image-only PDF skipped because OCR is disabled.")

    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or lowered_name.endswith(".docx")
    ):
        file_bytes = _download_file_bytes(drive_service, file_id)
        return _render_docx(file_bytes), "docx"

    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or lowered_name.endswith(".pptx")
    ):
        file_bytes = _download_file_bytes(drive_service, file_id)
        return _render_pptx(file_bytes), "pptx"

    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or lowered_name.endswith(".xlsx")
        or lowered_name.endswith(".xlsm")
    ):
        file_bytes = _download_file_bytes(drive_service, file_id)
        return _render_excel(file_bytes), "spreadsheet-file"

    if mime_type.startswith("text/") or any(lowered_name.endswith(suffix) for suffix in TEXT_FILE_SUFFIXES):
        file_bytes = _download_file_bytes(drive_service, file_id)
        return _render_text_like(file_bytes, lowered_name), "text-file"

    if _is_image_file(mime_type, lowered_name):
        raise ValueError("Image files are skipped because OCR is disabled.")

    raise ValueError(f"Unsupported file type: {mime_type}")


def _export_google_workspace_text(drive_service: Any, file_id: str) -> str:
    request = drive_service.files().export_media(fileId=file_id, mimeType="text/plain")
    handle = io.BytesIO()
    downloader = MediaIoBaseDownload(handle, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return handle.getvalue().decode("utf-8", errors="ignore")


def _download_file_bytes(drive_service: Any, file_id: str) -> bytes:
    request = drive_service.files().get_media(fileId=file_id)
    handle = io.BytesIO()
    downloader = MediaIoBaseDownload(handle, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return handle.getvalue()


def _render_google_sheet(sheets_service: Any, spreadsheet_id: str) -> str:
    spreadsheet = sheets_service.spreadsheets().get(
        spreadsheetId=spreadsheet_id,
        fields=SHEET_META_FIELDS,
    ).execute()

    title = spreadsheet.get("properties", {}).get("title", "Untitled spreadsheet")
    parts = [f"Spreadsheet: {title}"]
    sheet_names = [
        sheet.get("properties", {}).get("title", "Sheet")
        for sheet in spreadsheet.get("sheets", [])
    ]

    if not sheet_names:
        return "\n".join(parts)

    values_payload = sheets_service.spreadsheets().values().batchGet(
        spreadsheetId=spreadsheet_id,
        ranges=[_sheet_preview_range(sheet_name) for sheet_name in sheet_names],
        majorDimension="ROWS",
        valueRenderOption="FORMATTED_VALUE",
    ).execute()
    value_ranges = values_payload.get("valueRanges", [])

    for index, tab_name in enumerate(sheet_names):
        rows = value_ranges[index].get("values", []) if index < len(value_ranges) else []
        if not any(any(cell.strip() for cell in row) for row in rows):
            continue

        headers = rows[0] if rows else []
        data_rows = rows[1:] if len(rows) > 1 else []
        parts.append(f"\nTab: {tab_name}")

        if headers:
            header_text = " | ".join(header or f"Column {idx + 1}" for idx, header in enumerate(headers))
            parts.append(f"Columns: {header_text}")

        for row_number, row in enumerate(data_rows[:400], start=1):
            if headers:
                pairs = []
                for idx, cell in enumerate(row):
                    if not cell.strip():
                        continue
                    header = headers[idx] if idx < len(headers) and headers[idx] else f"Column {idx + 1}"
                    pairs.append(f"{header}: {cell}")
                if pairs:
                    parts.append(f"Row {row_number}: {' | '.join(pairs)}")
            else:
                flattened = " | ".join(cell for cell in row if cell.strip())
                if flattened:
                    parts.append(f"Row {row_number}: {flattened}")

    return "\n".join(parts)


def _sheet_preview_range(sheet_name: str) -> str:
    escaped_name = sheet_name.replace("'", "''")
    return f"'{escaped_name}'!1:{MAX_SHEET_ROWS}"


def _render_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"Page {index}\n{text}")
    return "\n\n".join(pages)


def _render_docx(file_bytes: bytes) -> str:
    document = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(paragraphs)


def _render_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    slides = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            stripped = text.strip()
            if stripped:
                texts.append(stripped)
        if texts:
            slides.append(f"Slide {slide_index}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _render_excel(file_bytes: bytes) -> str:
    workbook = pd.ExcelFile(io.BytesIO(file_bytes))
    sections = []

    for sheet_name in workbook.sheet_names:
        frame = workbook.parse(sheet_name=sheet_name).fillna("")
        if frame.empty:
            continue

        frame.columns = [str(column) if str(column).strip() else f"Column {idx + 1}" for idx, column in enumerate(frame.columns)]
        sections.append(f"Sheet: {sheet_name}")
        sections.append("Columns: " + " | ".join(frame.columns))

        for row_number, row in enumerate(frame.astype(str).to_dict(orient="records")[:400], start=1):
            pairs = [f"{column}: {value}" for column, value in row.items() if value.strip()]
            if pairs:
                sections.append(f"Row {row_number}: {' | '.join(pairs)}")

        sections.append("")

    return "\n".join(sections)


def _render_text_like(file_bytes: bytes, lowered_name: str) -> str:
    decoded = _decode_bytes(file_bytes)

    if lowered_name.endswith(".json"):
        parsed = json.loads(decoded)
        return json.dumps(parsed, indent=2, ensure_ascii=False)

    if lowered_name.endswith(".csv"):
        return _render_delimited_text(decoded, delimiter=",")

    if lowered_name.endswith(".tsv"):
        return _render_delimited_text(decoded, delimiter="\t")

    return decoded


def _render_delimited_text(raw_text: str, delimiter: str) -> str:
    reader = csv.reader(io.StringIO(raw_text), delimiter=delimiter)
    rows = list(reader)
    if not rows:
        return ""

    headers = rows[0]
    parts = ["Columns: " + " | ".join(headers)]
    for row_number, row in enumerate(rows[1:401], start=1):
        pairs = []
        for idx, value in enumerate(row):
            if not value.strip():
                continue
            column = headers[idx] if idx < len(headers) and headers[idx] else f"Column {idx + 1}"
            pairs.append(f"{column}: {value}")
        if pairs:
            parts.append(f"Row {row_number}: {' | '.join(pairs)}")
    return "\n".join(parts)


def _decode_bytes(file_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def _normalize_extracted_text(text: str) -> str:
    if not text:
        return ""

    lines = []
    for line in text.replace("\x00", " ").splitlines():
        collapsed = re.sub(r"\s+", " ", line).strip()
        if collapsed:
            lines.append(collapsed)

    return "\n".join(lines).strip()


def _needs_pdf_ocr(text: str) -> bool:
    compact = re.sub(r"\s+", "", text)
    return len(compact) < 80


def _is_image_file(mime_type: str, lowered_name: str) -> bool:
    if mime_type.startswith("image/") and mime_type != "image/svg+xml":
        return True
    return any(lowered_name.endswith(suffix) for suffix in IMAGE_FILE_SUFFIXES)
